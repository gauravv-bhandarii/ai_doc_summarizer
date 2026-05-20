import os
os.environ["TRANSFORMERS_NO_TF"] = "1"

from fastapi import FastAPI, UploadFile, File, Form, Depends, HTTPException, status
from fastapi.middleware.cors import CORSMiddleware
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from transformers import pipeline
from pydantic import BaseModel
from sqlalchemy.orm import Session
import pandas as pd
import io
import re

import models
from database import engine, get_db
import auth

models.Base.metadata.create_all(bind=engine)

app = FastAPI()

# Allow frontend to connect
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ❗ DO NOT load model here
summarizer = None
qa_pipeline = None

def extract_text(file, filename, max_chars=100000):
    text = ""
    ext = filename.split(".")[-1].lower() if "." in filename else ""

    try:
        if ext == "pdf":
            reader = PdfReader(file.file)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted: text += extracted + "\n"
                if len(text) > max_chars: break
        elif ext in ["docx", "doc"]:
            doc = docx.Document(file.file)
            for para in doc.paragraphs:
                text += para.text + "\n"
                if len(text) > max_chars: break
        elif ext == "pptx":
            ppt = Presentation(file.file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                if len(text) > max_chars: break
        elif ext in ["xlsx", "csv"]:
            contents = file.file.read()
            if ext == "csv":
                df = pd.read_csv(io.BytesIO(contents))
            else:
                df = pd.read_excel(io.BytesIO(contents))
            text = df.head(50).to_string()
        else:
            # Fallback for txt or others
            text = file.file.read()[:max_chars].decode("utf-8", errors="ignore")
    except Exception as e:
        print(f"Extraction error: {e}")
    
    # Reset file pointer if needed by multiple reads (though we don't right now)
    file.file.seek(0)
    return text.strip()

def find_best_context(question: str, text: str, max_context_len: int = 2500) -> str:
    paragraphs = [p.strip() for p in text.split('\n\n') if len(p.strip()) > 50]
    if not paragraphs:
        paragraphs = [p.strip() for p in text.split('\n') if len(p.strip()) > 50]
        if not paragraphs: return text[:max_context_len]

    q_words = set(re.findall(r'\w+', question.lower()))
    stopwords = {"what", "is", "the", "in", "of", "and", "a", "to", "for", "on", "are", "how", "why"}
    q_keywords = q_words - stopwords
    
    if not q_keywords: return text[:max_context_len]

    scored_paragraphs = []
    for p in paragraphs:
        p_words = set(re.findall(r'\w+', p.lower()))
        score = len(q_keywords.intersection(p_words))
        scored_paragraphs.append((score, p))
        
    scored_paragraphs.sort(key=lambda x: x[0], reverse=True)
    
    best_context = ""
    for score, p in scored_paragraphs:
        if len(best_context) + len(p) < max_context_len:
            best_context += p + "\n\n"
        else:
            break
            
    if not best_context: return text[:max_context_len]
    return best_context.strip()


@app.post("/summarize/")
async def summarize(file: UploadFile = File(...)):
    global summarizer

    # ✅ Load model only when needed
    if summarizer is None:
        try:
            summarizer = pipeline(
                "summarization",
                model="facebook/bart-large-cnn"
            )
        except Exception as e:
            # Log and return a clear error to the client
            print(f"Model load error: {e}")
            raise HTTPException(status_code=500, detail="Failed to load summarization model. Check server logs and internet connection.")

    text = extract_text(file, file.filename, max_chars=3000)

    # prevent overload
    text_to_summarize = text[:2000]

    if len(text_to_summarize) < 50:
        return {"summary": "Document is too short or could not extract readable text."}

    try:
        result = summarizer(
            text_to_summarize,
            max_length=150,
            min_length=50,
            do_sample=False
        )
    except Exception as e:
        print(f"Summarization error: {e}")
        raise HTTPException(status_code=500, detail="Summarization failed. See server logs for details.")

    return {"summary": result[0]["summary_text"]}


@app.post("/chat/")
async def chat_with_doc(question: str = Form(...), file: UploadFile = File(...)):
    global qa_pipeline

    if qa_pipeline is None:
        qa_pipeline = pipeline("question-answering", model="distilbert-base-cased-distilled-squad")

    text = extract_text(file, file.filename, max_chars=50000)
    
    # QA models have context limits, distilbert handles ~512 tokens.
    context = find_best_context(question, text, max_context_len=2500)

    if not context:
        return {"answer": "No readable text found in document to answer from."}

    result = qa_pipeline(question=question, context=context)
    
    # If confidence is extremely low, provide a disclaimer
    answer = result["answer"]
    score = result.get("score", 1.0)
    
    if score < 0.05:
        answer = f"I'm not completely sure, but it might be: {answer}"

    return {"answer": answer}

@app.post("/analyze_data/")
async def analyze_data(file: UploadFile = File(...)):
    filename = file.filename
    ext = filename.split(".")[-1].lower() if "." in filename else ""
    
    if ext not in ["csv", "xlsx"]:
        return {"error": "Unsupported file format for analysis"}
        
    try:
        contents = await file.read()
        if ext == "csv":
            df = pd.read_csv(io.BytesIO(contents))
        else:
            df = pd.read_excel(io.BytesIO(contents))
            
        num_cols = df.select_dtypes(include=['number']).columns.tolist()
        if not num_cols:
            return {"error": "No numerical data found for analysis"}
            
        total_rows = len(df)
        cat_cols = df.select_dtypes(include=['object', 'string']).columns.tolist()
        x_col = cat_cols[0] if cat_cols else df.columns[0]
        y_col = num_cols[0] if num_cols else None
        
        if y_col:
            sample_df = df.head(20).fillna(0)
            graph_data = []
            for _, row in sample_df.iterrows():
                graph_data.append({
                    "name": str(row[x_col])[:15],
                    "value": float(row[y_col])
                })
                
            metrics = {
                "Total Rows": total_rows,
                "Columns": len(df.columns),
                f"Avg {y_col}": float(df[y_col].mean()) if not df[y_col].empty else 0,
                f"Max {y_col}": float(df[y_col].max()) if not df[y_col].empty else 0
            }
            
            return {
                "type": "analytics",
                "graphData": graph_data,
                "metrics": metrics,
                "x_col": str(x_col),
                "y_col": str(y_col)
            }
    except Exception as e:
        return {"error": str(e)}
        
    return {"error": "Analysis failed"}

# --- AUTH & HISTORY ROUTES ---

class UserCreate(BaseModel):
    username: str
    password: str

class Token(BaseModel):
    access_token: str
    token_type: str

class HistoryCreate(BaseModel):
    filename: str
    file_size: str
    summary: str

@app.post("/register", response_model=Token)
def register(user: UserCreate, db: Session = Depends(get_db)):
    db_user = db.query(models.User).filter(models.User.username == user.username).first()
    if db_user:
        raise HTTPException(status_code=400, detail="Username already registered")
    hashed_password = auth.get_password_hash(user.password)
    new_user = models.User(username=user.username, hashed_password=hashed_password)
    db.add(new_user)
    db.commit()
    db.refresh(new_user)
    access_token = auth.create_access_token(data={"sub": new_user.username})
    return {"access_token": access_token, "token_type": "bearer"}

@app.post("/login", response_model=Token)
def login(user: UserCreate, db: Session = Depends(get_db)):
    db_user = db.query(models.User).filter(models.User.username == user.username).first()
    if not db_user or not auth.verify_password(user.password, db_user.hashed_password):
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Incorrect username or password",
            headers={"WWW-Authenticate": "Bearer"},
        )
    access_token = auth.create_access_token(data={"sub": db_user.username})
    return {"access_token": access_token, "token_type": "bearer"}

@app.post("/history/")
def create_history(history: HistoryCreate, username: str = Depends(auth.get_current_user), db: Session = Depends(get_db)):
    user = db.query(models.User).filter(models.User.username == username).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    db_history = models.History(
        user_id=user.id,
        filename=history.filename,
        file_size=history.file_size,
        summary=history.summary
    )
    db.add(db_history)
    db.commit()
    return {"status": "success"}

@app.get("/history/")
def get_history(username: str = Depends(auth.get_current_user), db: Session = Depends(get_db)):
    user = db.query(models.User).filter(models.User.username == username).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    histories = db.query(models.History).filter(models.History.user_id == user.id).order_by(models.History.created_at.desc()).all()
    # Format for frontend
    result = []
    for h in histories:
        result.append({
            "name": h.filename,
            "size": h.file_size,
            "summary": h.summary,
            "ts": h.created_at.isoformat()
        })
    return result